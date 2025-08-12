from flask import Flask, request, Response
from flasgger import Swagger
from flask_cors import CORS
import os
import pytesseract
from PIL import Image
import pdfplumber
import docx
import json
import csv
import io
import pandas as pd
from pptx import Presentation

# Optional: Set Tesseract path on Windows if needed
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

app = Flask(__name__)
Swagger(app)
CORS(app)

# ✅ Root route so Render home page shows something
@app.route("/")
def home():
    return Response("API is running ✅ — Go to /apidocs for Swagger UI", mimetype="text/plain")

@app.route("/extract", methods=["POST"])
def extract_file_content():
    """
    Upload a file and get plain text content
    ---
    consumes:
      - multipart/form-data
    parameters:
      - name: file
        in: formData
        type: file
        required: true
        description: Upload any file (.txt, .pdf, .jpg, .png, .docx, .pptx, .json, .csv, .xlsx)
    responses:
      200:
        description: Plain text content of the file
    """
    file = request.files.get("file")
    if not file:
        return Response("No file uploaded", status=400)

    filename = file.filename.lower()
    text = ""

    try:
        # TXT
        if filename.endswith(".txt"):
            text = file.read().decode("utf-8", errors="ignore")

        # PDF
        elif filename.endswith(".pdf"):
            with pdfplumber.open(file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"

        # Images (OCR)
        elif filename.endswith((".jpg", ".jpeg", ".png", ".bmp", ".tiff")):
            image = Image.open(file)
            text = pytesseract.image_to_string(image)

        # DOCX
        elif filename.endswith(".docx"):
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"

        # PPTX
        elif filename.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"

        # JSON
        elif filename.endswith(".json"):
            data = json.load(file)
            words = []
            def extract_from_json(obj):
                if isinstance(obj, dict):
                    for k, v in obj.items():
                        if k.lower() == "text" and isinstance(v, str):
                            words.append(v)
                        else:
                            extract_from_json(v)
                elif isinstance(obj, list):
                    for item in obj:
                        extract_from_json(item)

            extract_from_json(data)
            text = " ".join(words) if words else json.dumps(data, indent=2)

        # CSV
        elif filename.endswith(".csv"):
            reader = csv.reader(io.StringIO(file.read().decode("utf-8", errors="ignore")))
            for row in reader:
                text += " ".join(row) + "\n"

        # Excel
        elif filename.endswith((".xls", ".xlsx")):
            df = pd.read_excel(file)
            text = df.to_string(index=False)

        else:
            return Response(f"Unsupported file type: {filename}", status=400)

        return Response(text.strip(), mimetype="text/plain")

    except Exception as e:
        return Response(f"Error processing file: {str(e)}", status=500)

if __name__ == "__main__":
    app.run(debug=True)
